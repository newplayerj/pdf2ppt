from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from typing import List, Dict, Any, Optional
import logging
import os
import re
from PIL import Image

logger = logging.getLogger(__name__)

class PresentationGenerator:
    def __init__(self, output_path: str = "output/presentation.pptx"):
        """Initialize presentation generator with slide layouts."""
        self.prs = Presentation()
        self.output_path = output_path
        
        # Create and store slide layouts
        self.title_slide_layout = self.prs.slide_layouts[0]  # Title Slide
        self.section_slide_layout = self.prs.slide_layouts[1]  # Section Header
        self.content_slide_layout = self.prs.slide_layouts[1]  # Content Slide
        self.figure_slide_layout = self.prs.slide_layouts[5]  # Blank Slide for figures
        
        # Set default dimensions
        self.content_width = Inches(9)
        self.content_height = Inches(5)
        
        # Initialize slide layouts
        self.title_slide_layout = self.prs.slide_layouts[0]  # Title slide
        self.section_slide_layout = self.prs.slide_layouts[1]  # Section header
        self.bullet_slide_layout = self.prs.slide_layouts[1]   # Content with bullets
        self.figure_slide_layout = self.prs.slide_layouts[5]   # Figure slide
        
        # Set dimensions and styling
        self.left_margin = Inches(1)
        self.top_margin = Inches(1)
        self.width = Inches(8)
        self.height = Inches(5.5)
        
        # Text styling
        self.title_font_size = Pt(40)
        self.section_font_size = Pt(36)
        self.body_font_size = Pt(20)
        self.bullet_font_size = Pt(18)

    def _format_text_frame(self, text_frame, font_size=None):
        """Apply consistent formatting to text frame."""
        text_frame.word_wrap = True
        text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
        
        for paragraph in text_frame.paragraphs:
            paragraph.font.name = 'Arial'
            if font_size:
                paragraph.font.size = font_size
            paragraph.space_after = Pt(12)
            paragraph.space_before = Pt(6)

    def _add_title_slide(self, title: str, subtitle: str = None):
        """Create main title slide."""
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text_frame(title_shape.text_frame, self.title_font_size)
        
        # Add subtitle if provided
        if subtitle:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle
            self._format_text_frame(subtitle_shape.text_frame, self.section_font_size)

    def _add_section_slide(self, section_title: str):
        """Create section header slide."""
        slide = self.prs.slides.add_slide(self.section_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = section_title
        self._format_text_frame(title_shape.text_frame, self.section_font_size)

    def _add_content_slide(self, title: str, points: List[str]):
        """Create content slide with bullet points."""
        slide = self.prs.slides.add_slide(self.bullet_slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text_frame(title_shape.text_frame, self.section_font_size)
        
        # Add bullet points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for i, point in enumerate(points):
            # Clean the point text
            clean_point = point.strip().lstrip('•').strip()
            
            # Add paragraph
            if i == 0:
                p = tf.paragraphs[0]  # Use existing first paragraph
            else:
                p = tf.add_paragraph()
            
            p.text = clean_point
            p.level = 0
            
            # Format paragraph
            font = p.font
            font.size = self.bullet_font_size
            font.name = 'Arial'

    def _add_figure_description_slide(self, title: str, description: str):
        """Create a slide describing a figure's key points."""
        slide = self.prs.slides.add_slide(self.bullet_slide_layout)
        
        # Add title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text_frame(title_shape.text_frame, self.section_font_size)
        
        # Add description as bullet points
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        # Split description into bullet points
        points = self._extract_description_points(description)
        for i, point in enumerate(points):
            if i == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.text = point
            p.level = 0
            
            # Format paragraph
            font = p.font
            font.size = self.bullet_font_size
            font.name = 'Arial'

    def _extract_description_points(self, description: str) -> List[str]:
        """Extract key points from figure description."""
        # Split into sentences and clean up
        sentences = [s.strip() for s in description.split('.') if s.strip()]
        
        # Filter for relevant points about the figure
        points = []
        for sentence in sentences:
            # Look for sentences that likely describe the figure
            if any(keyword in sentence.lower() for keyword in 
                   ['show', 'illustrate', 'demonstrate', 'present', 'display', 
                    'indicate', 'reveal', 'depict', 'represent', 'compare']):
                points.append(sentence + '.')
        
        # If no points found, use all sentences
        if not points:
            points = [s + '.' for s in sentences]
        
        return points

    def _extract_figure_description(self, text: str, figure_ref: str) -> str:
        """Extract relevant description for a figure from the text."""
        sentences = text.split('. ')
        relevant_sentences = []
        
        # Find the figure reference
        for i, sentence in enumerate(sentences):
            if figure_ref.lower() in sentence.lower():
                # Include the reference sentence and the next 2 sentences for context
                relevant_sentences.extend(sentences[i:i+3])
        
        return '. '.join(relevant_sentences) if relevant_sentences else text

    def _match_figure(self, figure_info: Dict[str, Any], figures: List[str]) -> str:
        """
        Improved figure matching logic for dictionary figure information.
        Returns the best matching figure path or None if no match found.
        """
        try:
            # Extract figure reference from the dictionary
            fig_ref = figure_info.get('reference', '')
            
            # Clean up the reference
            ref_num = ''.join(filter(str.isdigit, fig_ref.split(':')[0]))
            
            if not ref_num:
                return None
                
            # Try different matching patterns
            patterns = [
                f"figure_{ref_num}",  # figure_1
                f"fig_{ref_num}",     # fig_1
                f"fig{ref_num}",      # fig1
                f"figure{ref_num}",   # figure1
                f"_{ref_num}_",       # _1_
                f"_{ref_num}.",       # _1.
            ]
            
            for pattern in patterns:
                matches = [
                    f for f in figures 
                    if pattern.lower() in f.lower()
                ]
                if matches:
                    logger.debug(f"Matched figure {fig_ref} to {matches[0]}")
                    return matches[0]
                    
            logger.warning(f"No matching figure found for {fig_ref}")
            return None
            
        except Exception as e:
            logger.error(f"Error matching figure: {str(e)}")
            return None

    def generate(self, content: Dict[str, Any], figures: List[str]) -> None:
        """Generate the presentation with improved figure handling."""
        try:
            # Add title slide
            title_slide = self.prs.slides.add_slide(self.title_slide_layout)
            title_slide.shapes.title.text = content['title']
            
            # Process each section
            for section in content['sections']:
                if section['title'] in ['References', 'Acknowledgements']:
                    continue
                    
                # Add section title slide
                self._add_section_slide(section['title'])
                
                # Add overview slide if available
                if 'overview' in section:
                    self._add_content_slide(
                        f"{section['title']} Overview",
                        [section['overview']]
                    )
                
                # Process content
                if 'content' in section:
                    for item in section['content']:
                        # Add content slides
                        if 'key_points' in item:
                            points = []
                            for point in item['key_points']:
                                # Format each point with its evidence
                                point_text = f"• {point['argument']}"
                                if point.get('evidence'):
                                    point_text += f"\n  - Evidence: {point['evidence']}"
                                if point.get('implications'):
                                    point_text += f"\n  - Impact: {point['implications']}"
                                points.append(point_text)
                            
                            self._add_content_slide(
                                item.get('subtitle', section['title']), 
                                points
                            )
                        
                        # Add figure slides
                        if 'figures' in item:
                            for figure_info in item['figures']:
                                matched_figure = self._match_figure(figure_info, figures)
                                
                                if matched_figure:
                                    # Create comprehensive description
                                    description = (
                                        f"{figure_info.get('description', '')}\n\n"
                                        f"Technical Details: {figure_info.get('technical_content', '')}\n"
                                        f"Results: {figure_info.get('results', '')}"
                                    )
                                    
                                    self._add_figure_slide(
                                        matched_figure,
                                        item.get('subtitle', section['title']),
                                        description
                                    )
                                else:
                                    logger.warning(f"Figure not found: {figure_info.get('reference')}")
            
            self.prs.save(self.output_path)
            logger.info(f"Presentation saved to {self.output_path}")
            
        except Exception as e:
            logger.error(f"Error generating presentation: {str(e)}")
            raise

    def _add_outline_slide(self, sections: List[str]):
        """Add an outline slide."""
        slide = self.prs.slides.add_slide(self.bullet_slide_layout)
        
        title_shape = slide.shapes.title
        title_shape.text = "Outline"
        
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for section in sections:
            p = tf.add_paragraph()
            p.text = section
            p.level = 0
            p.font.size = self.bullet_font_size
            p.font.name = 'Arial'

    def _extract_main_points(self, text: str) -> List[str]:
        """Extract main points from paragraph text."""
        # Split into sentences
        sentences = text.split('. ')
        # Filter and clean up points, removing any existing bullet points
        points = [
            s.strip().lstrip('•').strip() + '.' 
            for s in sentences 
            if len(s.strip()) > 20
        ]
        return points

    def _add_paper_title_slide(self, title: str):
        """Create the paper title slide with authors."""
        slide = self.prs.slides.add_slide(self.title_slide_layout)
        
        # Add paper title
        title_shape = slide.shapes.title
        title_shape.text = title
        self._format_text_frame(title_shape.text_frame, Pt(36))
        
        # Add authors
        subtitle_shape = slide.placeholders[1]
        authors = [
            "Tianhao Li",
            "Sandesh Shetty",
            "Advaith Kamath",
            "Ajay Jaiswal",
            "Xiaoqian Jiang",
            "Ying Ding",
            "Yejin Kim"
        ]
        subtitle_shape.text = "\n".join(authors)
        self._format_text_frame(subtitle_shape.text_frame, Pt(24))

    def _extract_figure_number(self, fig_ref: str) -> Optional[str]:
        """Extract figure number from reference text."""
        patterns = [
            r'Figure\s*(\d+)',
            r'Fig\.*\s*(\d+)',
            r'Figure\s*S(\d+)',
            r'Fig\.*\s*S(\d+)'
        ]
        
        for pattern in patterns:
            match = re.search(pattern, fig_ref)
            if match:
                return match.group(1)
        return None

    def _add_figure_slide(self, figure_path: str, title: str, description: str = None):
        """Add a slide with a figure and optional description."""
        slide = self.prs.slides.add_slide(self.figure_slide_layout)  # Changed from picture_slide_layout
        
        # Add title
        shapes = slide.shapes
        title_shape = shapes.title
        title_shape.text = title
        
        try:
            # Get image dimensions
            with Image.open(figure_path) as img:
                img_width, img_height = img.size
            
            # Calculate aspect ratio
            aspect_ratio = img_width / img_height
            
            # Calculate dimensions maintaining aspect ratio
            if aspect_ratio > self.content_width/self.content_height:
                width = self.content_width
                height = width / aspect_ratio
            else:
                height = self.content_height
                width = height * aspect_ratio
            
            # Calculate center position
            left = (Inches(10) - width) / 2
            top = Inches(2)  # Below title
            
            # Add picture
            slide.shapes.add_picture(figure_path, left, top, width, height)
            
            # Add description if provided
            if description:
                desc_left = Inches(1)
                desc_top = top + height + Inches(0.2)
                desc_width = Inches(8)
                desc_height = Inches(1)
                
                textbox = slide.shapes.add_textbox(desc_left, desc_top, desc_width, desc_height)
                textbox.text_frame.text = description
                textbox.text_frame.paragraphs[0].font.size = Pt(12)
                
        except Exception as e:
            logger.error(f"Error adding figure {figure_path}: {str(e)}")
            self._add_figure_placeholder(slide, "Figure could not be loaded")

    def _add_figure_placeholder(self, slide, message: str):
        """Add a placeholder shape when figure cannot be loaded."""
        left = Inches(2)
        top = Inches(2)
        width = Inches(6)
        height = Inches(4)
        
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, width, height
        )
        shape.fill.background()
        shape.line.color.rgb = RGBColor(128, 128, 128)
        
        # Add message in placeholder
        textbox = slide.shapes.add_textbox(
            left + Inches(0.5),
            top + Inches(1.5),
            width - Inches(1),
            height - Inches(3)
        )
        textbox.text_frame.text = message
        textbox.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER


